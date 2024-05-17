import mysql.connector
import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from apscheduler.schedulers.background import BackgroundScheduler
from datetime import datetime
import datetime
import time

mydb = mysql.connector.connect(
  host="localhost",
  user="root",
  password="",
  database="project"
)

mycurser = mydb.cursor()

sql = "select * from udata"
mycurser.execute(sql)

udata = mycurser.fetchall()

name = [i[1] for i in udata]
position = [i[2] for i in udata]
birthdate = [i[3] for i in udata]
graduation_year = [i[4] for i in udata]
gmail = [i[5] for i in udata]


def birthday_card():
  today = datetime.datetime.now().strftime("%Y-%m-%d")
  now_year = datetime.datetime.now().year
  for j in birthdate:
    n = birthdate.index(j)
    if today[6:10] == str(j)[6:10]:
        print(f"Today is {name[n]}'s birthday")
        graduation = graduation_year[n]
        prs = Presentation()

        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)


        slide_width = prs.slide_width
        slide_height = prs.slide_height


        bg_path = r"C:\projects\images\B2.png"
        left = top = 0
        width = slide_width
        height = slide_height
        slide.shapes.add_picture(bg_path, left, top, width, height)

        logo_path = r"images\B2.png"
        left = Inches(0.5)
        top = Inches(0.05)
        width =  Inches(3.0)
        height = Inches(0.5) 
        slide.shapes.add_picture(logo_path, left, top, width, height)

  
        left = (slide_width - Inches(5)) / 2
        top = Inches(2.7)  
        left = (slide_width - Inches(4)) / 2 + Inches(0.3)  
        width = Inches(3)
        height = Inches(2)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame


        

        text = "祝第X屆校友\n某某某現職\n生日快樂"
        text = text.replace("某某某", name[n])

        text_frame.text = text.strip()  # 刪除前後空格

        for paragraph in text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(30)
            paragraph.font.size = Inches(0.45)
            paragraph.font.name = '標楷體'
            paragraph.alignment = PP_ALIGN.CENTER

        fill = textbox.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(253, 255, 246) 
        line = textbox.line
        line.visible = False  

        # 保存修改後的ppt檔案
        prs.save(f"{name[n]}'s birthday_wish.pptx")

             
scheduler = BackgroundScheduler(timezone="Asia/Shanghai")

scheduler.add_job(birthday_card, 'cron', day_of_week='0-6', hour=4, minute=12) #生成賀卡時間

scheduler.start()

print('Schedule started ...')

while True: #刷新資料庫取得最新資料
    time.sleep(10) # 暫停10秒鐘
    mydb = mysql.connector.connect(
    host="localhost",
    user="root",
    password="",
    database="project"
  )

    mycurser = mydb.cursor()

    sql = "select * from udata"
    mycurser.execute(sql)

    udata = mycurser.fetchall()

    name = [i[1] for i in udata]
    position = [i[2] for i in udata]
    birthdate = [i[3] for i in udata]
    graduation_year = [i[4] for i in udata]
    gmail = [i[5] for i in udata]
    print('資料庫已更新.....')
